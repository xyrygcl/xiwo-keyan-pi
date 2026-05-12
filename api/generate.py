ifrom fastapi import FastAPI, File, UploadFile
from fastapi.responses import FileResponse
import os
import json
import subprocess
from pathlib import Path
import tempfile
import shutil

app = FastAPI()

# === PPT转视频的核心代码（和之前一样） ===
def probe_duration(audio_path):
    try:
        result = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "json", str(audio_path)],
            capture_output=True, text=True, check=True
        )
        return float(json.loads(result.stdout)["duration"])
    except:
        return 3.0

def escape_text(text):
    return text.replace("'", "'\\\\''").replace(":", "\\:")

def split_text(text, max_chars):
    if len(text) <= max_chars:
        return [text]
    lines = []
    current = ""
    for char in text:
        if len(current) >= max_chars:
            lines.append(current)
            current = char
        else:
            current += char
    if current:
        lines.append(current)
    return lines

def split_into_sentences(text):
    import re
    sentences = re.split(r'([。！？])', text)
    result = []
    for i in range(0, len(sentences)-1, 2):
        if sentences[i].strip():
            result.append(sentences[i] + (sentences[i+1] if i+1 < len(sentences) else ''))
    if len(sentences) % 2 == 1 and sentences[-1].strip():
        result.append(sentences[-1])
    return [s.strip() for s in result if s.strip()]

def choose_layout(slide_data):
    bullets = slide_data.get('bullets', [])
    paragraphs = slide_data.get('paragraphs', [])
    if len(bullets) >= 3:
        return 'three_column'
    elif len(paragraphs) > 0 and len(paragraphs[0]) > 100:
        return 'three_column'
    elif len(bullets) == 2 or len(paragraphs) == 2:
        return 'two_column'
    else:
        return 'single_column'

def generate_base_layout(duration, slide_num):
    filters = []
    filters.append(f"color=c=#f5f5f5:s=1920x1080:d={duration}[bg]")
    filters.append("[bg]drawbox=x=100:y=80:w=180:h=50:color=#ff9800:t=fill[bg1]")
    filters.append(
        "[bg1]drawtext=text='知识要点':"
        "fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf:fontsize=40:fontcolor=white:"
        "x=140:y=92[v0]"
    )
    return filters, "v0", 1

def layout_three_column(filters, current, layer_num, slide_data, slide_num):
    title = slide_data.get('title', f'第{slide_num}页')
    bullets = slide_data.get('bullets', [])
    paragraphs = slide_data.get('paragraphs', [])
    title_esc = escape_text(title)
    filters.append(
        f"[{current}]drawtext=text='{title_esc}':"
        f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf:fontsize=64:fontcolor=#1a1a1a:"
        f"x=100:y=160[v{layer_num}]"
    )
    current = f"v{layer_num}"
    layer_num += 1
    filters.append(
        f"[{current}]drawtext=text='重点看这几个维度。':"
        f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf:fontsize=32:fontcolor=#444444:"
        f"x=100:y=250[v{layer_num}]"
    )
    current = f"v{layer_num}"
    layer_num += 1
    if bullets and len(bullets) >= 3:
        contents = bullets[:3]
    elif paragraphs:
        sentences = split_into_sentences(paragraphs[0])
        if len(sentences) >= 3:
            per_card = len(sentences) // 3
            contents = [
                ''.join(sentences[:per_card]),
                ''.join(sentences[per_card:per_card*2]),
                ''.join(sentences[per_card*2:])
            ]
        elif len(sentences) == 2:
            contents = [sentences[0], sentences[1], ""]
        elif len(sentences) == 1:
            contents = [sentences[0], "", ""]
        else:
            contents = ["", "", ""]
    else:
        contents = ["", "", ""]
    for i in range(3):
        card_x = 100 + i * 580
        card_y = 350
        filters.append(
            f"[{current}]drawbox=x={card_x}:y={card_y}:w=520:h=600:color=white:t=fill[v{layer_num}]"
        )
        current = f"v{layer_num}"
        layer_num += 1
        filters.append(
            f"[{current}]drawbox=x={card_x}:y={card_y}:w=520:h=600:color=#e0e0e0:t=2[v{layer_num}]"
        )
        current = f"v{layer_num}"
        layer_num += 1
        filters.append(
            f"[{current}]drawtext=text='0{i+1}':"
            f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf:fontsize=48:fontcolor=#2196f3:"
            f"x={card_x+30}:y={card_y+30}[v{layer_num}]"
        )
        current = f"v{layer_num}"
        layer_num += 1
        if i < len(contents) and contents[i]:
            lines = split_text(contents[i][:150], 12)
            y = card_y + 110
            for line in lines[:12]:
                line_esc = escape_text(line)
                filters.append(
                    f"[{current}]drawtext=text='{line_esc}':"
                    f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf:fontsize=32:fontcolor=#444444:"
                    f"x={card_x+30}:y={y}[v{layer_num}]"
                )
                current = f"v{layer_num}"
                layer_num += 1
                y += 45
    return filters, current, layer_num

def layout_two_column(filters, current, layer_num, slide_data, slide_num):
    title = slide_data.get('title', f'第{slide_num}页')
    bullets = slide_data.get('bullets', [])
    paragraphs = slide_data.get('paragraphs', [])
    title_esc = escape_text(title)
    filters.append(
        f"[{current}]drawtext=text='{title_esc}':"
        f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf:fontsize=64:fontcolor=#1a1a1a:"
        f"x=100:y=160[v{layer_num}]"
    )
    current = f"v{layer_num}"
    layer_num += 1
    if bullets and len(bullets) >= 2:
        contents = bullets[:2]
    elif paragraphs and len(paragraphs) >= 2:
        contents = paragraphs[:2]
    elif paragraphs:
        para = paragraphs[0]
        mid = len(para) // 2
        contents = [para[:mid], para[mid:]]
    else:
        contents = ["", ""]
    for i in range(2):
        card_x = 150 + i * 900
        card_y = 300
        filters.append(
            f"[{current}]drawbox=x={card_x}:y={card_y}:w=820:h=650:color=white:t=fill[v{layer_num}]"
        )
        current = f"v{layer_num}"
        layer_num += 1
        filters.append(
            f"[{current}]drawbox=x={card_x}:y={card_y}:w=820:h=650:color=#e0e0e0:t=2[v{layer_num}]"
        )
        current = f"v{layer_num}"
        layer_num += 1
        if i < len(contents) and contents[i]:
            lines = split_text(contents[i][:200], 18)
            y = card_y + 50
            for line in lines[:14]:
                line_esc = escape_text(line)
                filters.append(
                    f"[{current}]drawtext=text='{line_esc}':"
                    f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf:fontsize=28:fontcolor=#444444:"
                    f"x={card_x+40}:y={y}[v{layer_num}]"
                )
                current = f"v{layer_num}"
                layer_num += 1
                y += 45
    return filters, current, layer_num

def layout_single_column(filters, current, layer_num, slide_data, slide_num):
    title = slide_data.get('title', f'第{slide_num}页')
    paragraphs = slide_data.get('paragraphs', [])
    title_esc = escape_text(title)
    filters.append(
        f"[{current}]drawtext=text='{title_esc}':"
        f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf:fontsize=72:fontcolor=#1a1a1a:"
        f"x=(w-text_w)/2:y=200[v{layer_num}]"
    )
    current = f"v{layer_num}"
    layer_num += 1
    filters.append(
        f"[{current}]drawbox=x=200:y=350:w=1520:h=600:color=white:t=fill[v{layer_num}]"
    )
    current = f"v{layer_num}"
    layer_num += 1
    filters.append(
        f"[{current}]drawbox=x=200:y=350:w=1520:h=600:color=#e0e0e0:t=2[v{layer_num}]"
    )
    current = f"v{layer_num}"
    layer_num += 1
    if paragraphs:
        lines = split_text(paragraphs[0][:300], 24)
        y = 400
        for line in lines[:12]:
            line_esc = escape_text(line)
            filters.append(
                f"[{current}]drawtext=text='{line_esc}':"
                f"fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf:fontsize=32:fontcolor=#444444:"
                f"x=250:y={y}[v{layer_num}]"
            )
            current = f"v{layer_num}"
            layer_num += 1
            y += 50
    return filters, current, layer_num

def generate_slide(slide_data, output_path, duration, slide_num):
    layout_type = choose_layout(slide_data)
    filters, current, layer_num = generate_base_layout(duration, slide_num)
    if layout_type == 'three_column':
        filters, current, layer_num = layout_three_column(filters, current, layer_num, slide_data, slide_num)
    elif layout_type == 'two_column':
        filters, current, layer_num = layout_two_column(filters, current, layer_num, slide_data, slide_num)
    else:
        filters, current, layer_num = layout_single_column(filters, current, slide_data, slide_num)
    filter_complex = ";".join(filters)
    cmd = ["ffmpeg", "-y"]
    cmd.extend(["-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo"])
    cmd.extend([
        "-filter_complex", filter_complex,
        "-map", f"[{current}]", "-map", "0:a",
        "-c:v", "libx264", "-pix_fmt", "yuv420p",
        "-c:a", "aac", "-b:a", "192k",
        "-t", str(duration), "-shortest", str(output_path)
    ])
    subprocess.run(cmd, check=True, capture_output=True)
    return layout_type

def parse_pptx(ppt_path):
    from pptx import Presentation
    prs = Presentation(ppt_path)
    slides = []
    for slide in prs.slides:
        title = ""
        paragraphs = []
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                title = shape.text
            else:
                for p in shape.text_frame.paragraphs:
                    if p.level == 0:
                        bullets.append(p.text)
                    else:
                        paragraphs.append(p.text)
        slides.append({
            "title": title or f"第{len(slides)+1}页",
            "paragraphs": paragraphs,
            "bullets": bullets
        })
    return slides

def generate_video(slides, output_path):
    temp_dir = tempfile.mkdtemp()
    segments = []
    for i, slide in enumerate(slides, 1):
        duration = 5.0
        segment = Path(temp_dir) / f"seg_{i:03d}.mp4"
        generate_slide(slide, segment, duration, i)
        segments.append(segment)
    concat_file = Path(temp_dir) / "concat.txt"
    with open(concat_file, "w") as f:
        for seg in segments:
            f.write(f"file '{seg}'\n")
    subprocess.run([
        "ffmpeg", "-y", "-f", "concat", "-safe", "0",
        "-i", str(concat_file), "-c", "copy", str(output_path)
    ], check=True)
    shutil.rmtree(temp_dir)

@app.post("/api/generate")
async def generate(ppt: UploadFile = File(...)):
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        ppt_path = tmp_path / "input.pptx"
        with open(ppt_path, "wb") as f:
            shutil.copyfileobj(ppt.file, f)
        slides = parse_pptx(ppt_path)
        video_path = tmp_path / "output.mp4"
        generate_video(slides, video_path)
        return FileResponse(video_path, filename="微课视频.mp4", media_type="video/mp4")

@app.get("/")
async def root():
    return FileResponse("index.html")
